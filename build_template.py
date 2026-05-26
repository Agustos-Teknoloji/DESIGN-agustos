"""
Build agustos-deck-template.pptx — a reusable template demonstrating each
slide layout in the Ağustos design system.

Use this as a starting point for any new deck. Replace placeholder copy
on each slide; the type tokens, brand colors, and layout positions stay
constant.

Source of truth: PROJECTS/DESIGN-agustos/DESIGN.md
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Tokens (mirrors DESIGN.md) ---
PAPER     = RGBColor(0xfe, 0xfc, 0xf2)
INK       = RGBColor(0x1a, 0x1a, 0x1a)
INK_SOFT  = RGBColor(0x4a, 0x4a, 0x4a)
INK_FAINT = RGBColor(0x8a, 0x8a, 0x8a)
RULE      = RGBColor(0xe8, 0xe3, 0xd0)
BRAND     = RGBColor(0xcf, 0x14, 0x2a)

SERIF     = "Cambria"        # Newsreader fallback
SANS      = "Calibri"        # Source Sans 3 fallback
LOGOTYPE  = "Calibri Light"  # IBM Plex Sans Light fallback

def _resolve(paths):
    for p in paths:
        if os.path.exists(p):
            return p
    return paths[0]

SYMBOL_RED = _resolve([
    "/sessions/optimistic-cool-franklin/mnt/business/PROJECTS/DESIGN-agustos/laz-gunesi-amblem/png/laz-gunesi__red_2048px.png",
])
OUTPUT_DIR = "/sessions/optimistic-cool-franklin/mnt/outputs"
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "agustos-deck-template.pptx")

# --- Helpers (subset — mirror build_deck.py) ---

def set_bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, x, y, w, h, text, *, font=SERIF, size=16, bold=False,
             italic=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking=0, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if tracking:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(tracking * 100)))
    return tb

def add_paragraphs(slide, x, y, w, h, paragraphs, *, align=PP_ALIGN.LEFT,
                   anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", align)
        if "space_before" in para: p.space_before = Pt(para["space_before"])
        if "space_after" in para:  p.space_after = Pt(para["space_after"])
        if "line_spacing" in para: p.line_spacing = para["line_spacing"]
        for run in para["runs"]:
            r = p.add_run()
            r.text = run["text"]
            f = r.font
            f.name = run.get("font", SERIF)
            f.size = Pt(run.get("size", 16))
            f.bold = run.get("bold", False)
            f.italic = run.get("italic", False)
            f.color.rgb = run.get("color", INK)
            if "tracking" in run:
                rPr = r._r.get_or_add_rPr()
                rPr.set("spc", str(int(run["tracking"] * 100)))
    return tb

def add_rule(slide, x, y, w, h=0.025, color=BRAND):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def add_vrule(slide, x, y, h, w=0.025, color=BRAND):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(w), h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

# Token convenience wrappers
def eyebrow(slide, x, y, text, w=Inches(8)):
    return add_text(slide, x, y, w, Inches(0.3),
                    text.upper(), font=SANS, size=12, bold=True,
                    color=INK_FAINT, tracking=0.14)

def h1(slide, x, y, text, *, w=Inches(11.5), size=44, color=INK):
    return add_text(slide, x, y, w, Inches(1.6),
                    text, font=SERIF, size=size, color=color,
                    line_spacing=1.05, tracking=-0.024)

def deck_italic(slide, x, y, text, *, w=Inches(11.5), size=22):
    return add_text(slide, x, y, w, Inches(0.9),
                    text, font=SERIF, size=size, italic=True, color=INK_SOFT,
                    line_spacing=1.35, tracking=-0.012)

def body(slide, x, y, text, *, w=Inches(10.5), size=17):
    return add_text(slide, x, y, w, Inches(2),
                    text, font=SERIF, size=size, color=INK,
                    line_spacing=1.55)

def add_lockup(slide, x, y, *, height=0.42):
    sym_size = Inches(height)
    pic = slide.shapes.add_picture(SYMBOL_RED, x, y,
                                   width=sym_size, height=sym_size)
    word_size_pt = round((height / 1.4) * 72 * 1.25)
    gap = Inches(height * 0.5)
    add_text(slide, x + sym_size + gap, y, Inches(2.4), sym_size,
             "ağustos",
             font=LOGOTYPE, size=word_size_pt, color=BRAND,
             anchor=MSO_ANCHOR.MIDDLE, tracking=-0.02)
    return pic

def add_footer_mark(slide, *, page_meta="Section · Page title  ·  YEAR"):
    fy = SLIDE_H - Inches(0.55)
    sym_size = Inches(0.28)
    slide.shapes.add_picture(SYMBOL_RED, Inches(0.55), fy,
                             width=sym_size, height=sym_size)
    add_text(slide, Inches(0.55) + sym_size + Inches(0.16), fy,
             Inches(5), sym_size,
             "ağustos teknoloji", font=LOGOTYPE, size=11, color=INK_SOFT,
             tracking=-0.01, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(2.5) - Inches(0.55), fy,
             Inches(2.5), sym_size,
             page_meta, font=SANS, size=9, color=INK_FAINT,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, tracking=0.08)

# --- Presentation setup ---
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# Helper note: small grey caption explaining what each slide demonstrates.
def template_note(slide, text):
    add_text(slide, Inches(0.55), SLIDE_H - Inches(1.0),
             Inches(SLIDE_W.inches - 1.1), Inches(0.3),
             text, font=SANS, size=9, italic=True,
             color=INK_FAINT, align=PP_ALIGN.LEFT)

# ---------- Slide 1 — Cover note about the template itself ----------
def slide_cover():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Template · v1.0")
    add_text(s, Inches(0.9), Inches(1.25), Inches(11.5), Inches(2.2),
             "Ağustos slide template.",
             font=SERIF, size=64, color=INK,
             line_spacing=1.02, tracking=-0.024)
    add_text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.9),
             "Each slide that follows demonstrates one layout pattern. "
             "Replace the placeholder copy; keep the position, size, and color tokens.",
             font=SERIF, size=22, italic=True, color=INK_SOFT,
             line_spacing=1.35, tracking=-0.012)
    add_rule(s, Inches(0.9), Inches(4.7), Inches(2.0))
    # Token cheatsheet
    add_paragraphs(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(2.0), [
        {"runs": [
            {"text": "PAPER ", "font": SANS, "size": 10, "tracking": 0.14,
             "color": INK_FAINT, "bold": True},
            {"text": "  #fefcf2", "size": 14, "color": INK},
            {"text": "       INK ", "font": SANS, "size": 10, "tracking": 0.14,
             "color": INK_FAINT, "bold": True},
            {"text": "  #1a1a1a", "size": 14, "color": INK},
            {"text": "       BRAND ", "font": SANS, "size": 10, "tracking": 0.14,
             "color": INK_FAINT, "bold": True},
            {"text": "  #cf142a", "size": 14, "color": BRAND},
        ], "space_after": 4},
        {"runs": [
            {"text": "SERIF ", "font": SANS, "size": 10, "tracking": 0.14,
             "color": INK_FAINT, "bold": True},
            {"text": "  Newsreader · Cambria fallback", "size": 14, "color": INK},
        ], "space_after": 4},
        {"runs": [
            {"text": "SANS ", "font": SANS, "size": 10, "tracking": 0.14,
             "color": INK_FAINT, "bold": True},
            {"text": "  Source Sans 3 · Calibri fallback", "size": 14, "color": INK},
        ], "space_after": 4},
        {"runs": [
            {"text": "LOGOTYPE ", "font": SANS, "size": 10, "tracking": 0.14,
             "color": INK_FAINT, "bold": True},
            {"text": "  IBM Plex Sans Light · Calibri Light fallback", "size": 14,
             "color": INK},
        ]},
    ])
    add_lockup(s, Inches(0.9), Inches(6.85), height=0.4)
    return s

# ---------- Slide 2 — Title ----------
def slide_title():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Layout · Title")
    add_text(s, Inches(0.9), Inches(1.25), Inches(11.5), Inches(1.6),
             "Headline goes here.",
             font=SERIF, size=64, color=INK,
             line_spacing=1.02, tracking=-0.024)
    add_text(s, Inches(0.9), Inches(2.55), Inches(11), Inches(0.9),
             "Italic deck — the second sentence after the title.",
             font=SERIF, size=24, italic=True, color=INK_SOFT,
             line_spacing=1.3, tracking=-0.012)
    # placeholder for hero illustration (rule + caption)
    add_rule(s, Inches(0.9), Inches(3.85), Inches(11.5),
             color=RULE, h=0.012)
    add_text(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(1.6),
             "[ Optional hero illustration — wide horizontal band ]",
             font=SANS, size=11, italic=True, color=INK_FAINT,
             align=PP_ALIGN.CENTER)
    add_rule(s, Inches(0.9), Inches(5.85), Inches(2.0))
    add_paragraphs(s, Inches(0.9), Inches(6.05), Inches(8.5), Inches(0.9), [
        {"runs": [
            {"text": "PRESENTED BY ", "font": SANS, "size": 10, "tracking": 0.14,
             "color": INK_FAINT, "bold": True},
            {"text": "  Author Name", "size": 15, "bold": True},
            {"text": "  ·  Ağustos Teknoloji", "size": 15,
             "color": INK_SOFT, "italic": True},
        ]},
    ])
    add_lockup(s, Inches(0.9), Inches(6.85), height=0.4)
    add_text(s, SLIDE_W - Inches(2.5) - Inches(0.9), Inches(6.9),
             Inches(2.5), Inches(0.35),
             "YEAR", font=SANS, size=12, color=INK_FAINT,
             align=PP_ALIGN.RIGHT, tracking=0.08)
    return s

# ---------- Slide 3 — Section divider ----------
def slide_divider():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Layout · Section divider")
    sym_size = Inches(3.0)
    s.shapes.add_picture(SYMBOL_RED,
                         (SLIDE_W - sym_size) / 2, Inches(1.6),
                         width=sym_size, height=sym_size)
    add_text(s, Inches(1), Inches(5.05), Inches(11.3), Inches(1.5),
             "Why?", font=SERIF, size=80, color=INK,
             align=PP_ALIGN.CENTER, line_spacing=1.0, tracking=-0.024)
    add_text(s, Inches(1), Inches(6.65), Inches(11.3), Inches(0.5),
             "One-sentence subtitle in italic.",
             font=SERIF, size=17, italic=True, color=INK_SOFT,
             align=PP_ALIGN.CENTER, tracking=-0.006)
    return s

# ---------- Slide 4 — Content (full-width body) ----------
def slide_content_full():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Layout · Content (full width)")
    h1(s, Inches(0.9), Inches(1.15), "Headline: a complete sentence.")
    deck_italic(s, Inches(0.9), Inches(2.75),
                "An italic deck under the H1 — sets up the body.")
    add_rule(s, Inches(0.9), Inches(4.05), Inches(2.0))
    body(s, Inches(0.9), Inches(4.4),
         "Body copy in the serif. Use this layout when the slide is "
         "primarily prose — a definition, a position, a long quote. The "
         "italic deck under the H1 sets up the body; the brand-red rule "
         "separates them.",
         w=Inches(11.5), size=17)
    add_footer_mark(s)
    return s

# ---------- Slide 5 — Content (with image right) ----------
def slide_content_image_right():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Layout · Content with image")
    h1(s, Inches(0.9), Inches(1.15), "Headline next to evidence.")
    deck_italic(s, Inches(0.9), Inches(2.75),
                "Italic deck — fits one line at this width.",
                w=Inches(6.4))
    add_rule(s, Inches(0.9), Inches(4.05), Inches(2.0))
    add_paragraphs(s, Inches(0.9), Inches(4.35), Inches(6.4), Inches(2.5), [
        {"runs": [
            {"text": "— ", "size": 17, "color": BRAND, "bold": True},
            {"text": "First bullet point", "size": 17},
        ], "space_after": 4, "line_spacing": 1.4},
        {"runs": [
            {"text": "— ", "size": 17, "color": BRAND, "bold": True},
            {"text": "Second bullet point", "size": 17},
        ], "space_after": 4, "line_spacing": 1.4},
        {"runs": [
            {"text": "— ", "size": 17, "color": BRAND, "bold": True},
            {"text": "Third bullet point", "size": 17},
        ], "space_after": 4, "line_spacing": 1.4},
    ])
    # image placeholder
    img_x, img_y, img_w, img_h = Inches(7.6), Inches(3.0), Inches(5.0), Inches(3.4)
    ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_x, img_y, img_w, img_h)
    ph.fill.solid(); ph.fill.fore_color.rgb = RULE
    ph.line.color.rgb = INK_FAINT
    add_text(s, img_x, img_y + img_h/2 - Inches(0.15), img_w, Inches(0.3),
             "[ Image — 5.0\" × 3.4\" ]",
             font=SANS, size=12, italic=True, color=INK_FAINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, img_x, img_y + img_h + Inches(0.1), img_w, Inches(0.3),
             "Optional italic caption in faint ink.",
             font=SANS, size=10, italic=True, color=INK_FAINT,
             align=PP_ALIGN.CENTER)
    add_footer_mark(s)
    return s

# ---------- Slide 6 — Comparison (Before / After) ----------
def slide_comparison():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Layout · Two-column comparison")
    h1(s, Inches(0.9), Inches(1.15),
       "From one state to another, side by side.", size=40)
    add_rule(s, Inches(0.9), Inches(3.0), Inches(2.0))

    col_y = Inches(3.4)
    col_w = Inches(5.6)
    col_lx, col_rx = Inches(0.9), Inches(6.9)

    # Before column
    add_text(s, col_lx, col_y, col_w, Inches(0.4),
             "BEFORE", font=SANS, size=12, bold=True, color=INK_FAINT,
             tracking=0.14)
    add_paragraphs(s, col_lx, col_y + Inches(0.5), col_w, Inches(3), [
        {"runs": [
            {"text": "Attribute: ", "size": 18, "italic": True, "color": INK},
            {"text": "value", "size": 18, "color": INK_SOFT},
        ], "space_after": 8},
        {"runs": [
            {"text": "Attribute: ", "size": 18, "italic": True, "color": INK},
            {"text": "value", "size": 18, "color": INK_SOFT},
        ], "space_after": 8},
        {"runs": [
            {"text": "Attribute: ", "size": 18, "italic": True, "color": INK},
            {"text": "value", "size": 18, "color": INK_SOFT},
        ], "space_after": 8},
    ])

    # After column with brand-red accent rule
    add_vrule(s, col_rx - Inches(0.18), col_y + Inches(0.05),
              Inches(0.45), w=0.025)
    add_text(s, col_rx, col_y, col_w, Inches(0.4),
             "AFTER", font=SANS, size=12, bold=True, color=BRAND,
             tracking=0.14)
    add_paragraphs(s, col_rx, col_y + Inches(0.5), col_w, Inches(3), [
        {"runs": [
            {"text": "Attribute: ", "size": 18, "italic": True, "color": INK},
            {"text": "value", "size": 18, "color": INK},
        ], "space_after": 8},
        {"runs": [
            {"text": "Attribute: ", "size": 18, "italic": True, "color": INK},
            {"text": "value", "size": 18, "color": INK},
        ], "space_after": 8},
        {"runs": [
            {"text": "Attribute: ", "size": 18, "italic": True, "color": INK},
            {"text": "value", "size": 18, "color": INK},
        ], "space_after": 8},
    ])
    add_footer_mark(s)
    return s

# ---------- Slide 7 — Step (giant numeral + bullets) ----------
def slide_step():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Layout · Step / numbered phase")
    add_text(s, Inches(0.9), Inches(1.1), Inches(2.5), Inches(2.6),
             "1", font=SERIF, size=140, color=BRAND,
             line_spacing=1.0, tracking=-0.024)
    add_text(s, Inches(0.9), Inches(3.7), Inches(2.5), Inches(0.4),
             "STEP 1", font=SANS, size=12, bold=True, color=INK_FAINT,
             tracking=0.14)
    h1(s, Inches(4.5), Inches(1.15), "Step name.", w=Inches(8.5), size=42)
    add_text(s, Inches(4.5), Inches(2.95), Inches(8.5), Inches(1.0),
             "Italic deck describing what this step is.",
             font=SERIF, size=20, italic=True, color=INK_SOFT,
             line_spacing=1.35, tracking=-0.012)
    add_rule(s, Inches(4.5), Inches(4.15), Inches(2.0))
    add_paragraphs(s, Inches(4.5), Inches(4.4), Inches(8.5), Inches(2.5), [
        {"runs": [
            {"text": "— ", "size": 16, "color": BRAND, "bold": True},
            {"text": "First activity in this step", "size": 16},
        ], "space_after": 5, "line_spacing": 1.45},
        {"runs": [
            {"text": "— ", "size": 16, "color": BRAND, "bold": True},
            {"text": "Second activity in this step", "size": 16},
        ], "space_after": 5, "line_spacing": 1.45},
        {"runs": [
            {"text": "— ", "size": 16, "color": BRAND, "bold": True},
            {"text": "Third activity in this step", "size": 16},
        ], "space_after": 5, "line_spacing": 1.45},
    ])
    add_footer_mark(s)
    return s

# ---------- Slide 8 — Diagram (full-width image) ----------
def slide_diagram():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.55), "Layout · Diagram")
    add_text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.75),
             "Diagram or flow goes here.",
             font=SERIF, size=32, color=INK,
             line_spacing=1.1, tracking=-0.012)
    add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.85),
             "Italic deck framing what the diagram shows.",
             font=SERIF, size=17, italic=True, color=INK_SOFT,
             line_spacing=1.4, tracking=-0.006)
    # image placeholder
    img_x, img_y = Inches(2.55), Inches(2.15)
    img_w, img_h = Inches(8.225), Inches(4.7)
    ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_x, img_y, img_w, img_h)
    ph.fill.solid(); ph.fill.fore_color.rgb = RULE
    ph.line.color.rgb = INK_FAINT
    add_text(s, img_x, img_y + img_h/2 - Inches(0.2), img_w, Inches(0.4),
             "[ Diagram — 14:9.6 aspect, 1400×800 native ]",
             font=SANS, size=14, italic=True, color=INK_FAINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer_mark(s)
    return s

# ---------- Slide 9 — Closing ----------
def slide_close():
    s = prs.slides.add_slide(BLANK); set_bg(s)
    eyebrow(s, Inches(0.9), Inches(0.7), "Layout · Closing")
    sym_size = Inches(2.0)
    s.shapes.add_picture(SYMBOL_RED,
                         (SLIDE_W - sym_size) / 2, Inches(1.4),
                         width=sym_size, height=sym_size)
    add_text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(1.2),
             "Thank you.", font=SERIF, size=64, color=INK,
             align=PP_ALIGN.CENTER, line_spacing=1.05, tracking=-0.024)
    add_text(s, Inches(1), Inches(4.85), Inches(11.3), Inches(0.6),
             "One italic line — invitation, contact, or next step.",
             font=SERIF, size=20, italic=True, color=INK_SOFT,
             align=PP_ALIGN.CENTER, line_spacing=1.4, tracking=-0.012)
    add_lockup(s, Inches(0.55), Inches(6.85), height=0.34)
    add_text(s, SLIDE_W - Inches(2.5) - Inches(0.55), Inches(6.9),
             Inches(2.5), Inches(0.3),
             "Section · Page title  ·  YEAR",
             font=SANS, size=9, color=INK_FAINT,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, tracking=0.08)
    return s

# --- Build ---
slide_cover()
slide_title()
slide_divider()
slide_content_full()
slide_content_image_right()
slide_comparison()
slide_step()
slide_diagram()
slide_close()

prs.save(OUTPUT_FILE)
print(f"Wrote {OUTPUT_FILE}")
print(f"Slides: {len(prs.slides)}")
